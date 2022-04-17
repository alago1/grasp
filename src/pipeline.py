import io
from os import path
from typing import List, Tuple

from PIL.Image import Image
from pptx import Presentation
from pptx.util import Inches

from buffer import get_frame_data
from reference_producer import get_references
from text_detection import detect_text

#  pylint: disable=unspecified-encoding


def add_annotated_slide(prs: Presentation, screenshot: Image, annotation: str) -> None:
    """
    Given a presentation instance, an image screenshot, and an annotation string,
    adds the annotated screenshot to a new blank slide in the presentation.
    """

    blank_slide = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide)

    with io.BytesIO() as output:
        screenshot.save(output, format="GIF")
        slide.shapes.add_picture(output, Inches(0), Inches(0), width=prs.slide_width)
        slide.notes_slide.notes_text_frame.text = annotation


def build_output_path(output_folder: str, input_filepath: str, output_type: str) -> str:
    """
    Given the output_folder, input_filepath, and output_type,
    constructs the output_path with the the proper file type and function.
    """

    filename = path.split(input_filepath)[-1].split(".")[0]

    if output_type == "PPTX":
        return path.join(output_folder, f"slides_{filename}.pptx")

    return path.join(output_folder, f"references_{filename}.txt")


def export_txt(
    output_path: str, references: List[Tuple[int, List[Tuple[str, str]]]]
) -> None:
    """
    Given a list of references, exports the timestamp, keywords and
    reference to a text file in markdown format.
    """
    with open(output_path, "w") as output_file:
        output_file.write("# References:\n")

        seen_keywords = set()

        for timestamp, refs in references:
            filtered_refs = []
            for keyword, url in refs:
                if keyword in seen_keywords:
                    continue
                seen_keywords.add(keyword)
                filtered_refs.append((keyword, url))

            if filtered_refs:
                output_file.write(f"## CURRENT TIMESTAMP: {timestamp}\n")
                for item in filtered_refs:
                    output_file.write(f"{item}\n")


def process_references(
    filepath: str, output_folder: str, output_type: str, disable_exporting: bool = False
) -> List[Tuple[int, List[Tuple[str, str]]]]:
    """
    Given filepath to video file, process entire
    pipeline and returns list of timestamped pairs
    of noun-chunks, references.
    """

    prs = Presentation()
    references: List[Tuple[int, List[Tuple[str, str]]]] = []

    img_buffer = get_frame_data(filepath)
    for timestamp, img in img_buffer:
        print(f"CURRENT TIMESTAMP: {timestamp}")
        text = detect_text(img)
        refs = get_references(text)

        references.append((timestamp, refs))

        if not disable_exporting and output_type == "PPTX":
            annotation = "\n".join([f'"{kw}": {ref}' for kw, ref in refs])
            add_annotated_slide(prs, img, annotation)

    if disable_exporting:
        return references

    output_path = build_output_path(output_folder, filepath, output_type)

    if output_type == "PPTX":
        prs.save(output_path)
    else:
        export_txt(
            output_path,
            references,
        )

    return references
